"""Office documents generated at test time rather than committed as binaries.

Same reasoning as `tests/fixtures_pdf.py`: a committed binary rots, nobody can
read a diff of one, and a corrupted byte looks identical to an intentional
edit. Generating them keeps the input to every office test readable in the same
repository as the test.

Two of these do something python-docx and openpyxl cannot, by rewriting one
part inside the finished zip:

* **Tracked changes.** python-docx has no API for `w:ins`/`w:del`, and the Word
  handler reports their presence as a card fact. Injecting the element is the
  only way to have a document that carries one.
* **Cached formula values.** openpyxl never computes a formula and writes no
  cached `<v>`, so `data_only=True` returns `None` for every formula cell in
  any workbook it produced. A test of "represent shows the value" against such
  a workbook would be asserting that a blank is a blank. Injecting the `<v>`
  Excel itself would have written is what makes the value path real.

ODF is hand-built. There is no maintained writer, the format is a three-entry
zip, and building it here means the test knows exactly what it fed the reader.

python-docx, python-pptx and openpyxl are test-time writers here. Nothing under
`src/` imports them outside its one handler, and within `tests/` they are
confined to this module and its guard test.
"""

from __future__ import annotations

import io
import re
import zipfile
from collections.abc import Callable, Sequence

import docx
import openpyxl
import pptx
from pptx.util import Inches

#: ODF namespaces, only the ones these fixtures actually emit.
_OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
_TABLE_NS = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
_DRAW_NS = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"

_MANIFEST = (
    '<?xml version="1.0" encoding="UTF-8"?>'
    "<manifest:manifest xmlns:manifest="
    '"urn:oasis:names:tc:opendocument:xmlns:manifest:1.0" manifest:version="1.3"/>'
)

#: What `_cache_formula_values` writes as every formula's cached result. A
#: constant, not the formula's real answer: the point is that a cached value
#: EXISTS and is what the value path reports, not that this module can do
#: arithmetic.
CACHED_FORMULA_VALUE = 99


def _rewrite_part(data: bytes, name: str, edit: Callable[[bytes], bytes]) -> bytes:
    """The package with one part passed through `edit`, everything else copied.

    Entry ORDER is preserved, which matters more than it looks:
    `adapters/ooxml.office_mimetype` reads local headers in stored order, and an
    ODF package is only sniffable because `mimetype` is first. A naive rebuild
    that reordered entries would make the detection tests pass or fail for
    reasons that have nothing to do with detection.
    """
    source = zipfile.ZipFile(io.BytesIO(data))
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            blob = source.read(item.filename)
            if item.filename == name:
                blob = edit(blob)
            target.writestr(item, blob)
    return buffer.getvalue()


def _insert_tracked_change(xml: bytes) -> bytes:
    """Wrap the first run in a `w:ins`, the way Word records an insertion.

    python-docx has no API for this and the Word handler reports tracked
    changes as a card fact, so without this there is no document that exercises
    the true branch — and a fact that is only ever "no" tests nothing.
    """
    opening = re.search(rb"<w:r>", xml)
    if opening is None:
        return xml
    closing = xml.index(b"</w:r>", opening.end()) + len(b"</w:r>")
    run = xml[opening.start() : closing]
    inserted = (
        b'<w:ins w:id="901" w:author="Reviewer" w:date="2026-08-15T00:00:00Z">' + run + b"</w:ins>"
    )
    return xml[: opening.start()] + inserted + xml[closing:]


def docx_bytes(
    *,
    headings: Sequence[tuple[int, str, str]] = (
        (1, "Alpha", "The body of the alpha section."),
        (2, "Bravo", "The body of the bravo section."),
        (1, "Charlie", "The body of the charlie section."),
    ),
    table: Sequence[Sequence[str]] | None = (("region", "total"), ("north", "12")),
    comment: str | None = None,
    tracked: bool = False,
    preamble: str | None = None,
) -> bytes:
    """A Word document with headings, body paragraphs and optionally a table.

    The table is placed after the FIRST heading rather than at the end, so a
    handler that appends tables instead of reading document order produces
    visibly wrong text.

    `preamble` is body text before any heading — the case that proves a
    `LocatorMap` starting at offset 0 has something to attribute.
    """
    document = docx.Document()
    if preamble is not None:
        document.add_paragraph(preamble)
    for index, (level, title, body) in enumerate(headings):
        document.add_heading(title, level)
        paragraph = document.add_paragraph(body)
        if index == 0:
            if table is not None:
                built = document.add_table(rows=len(table), cols=len(table[0]))
                for row_index, row in enumerate(table):
                    for cell_index, value in enumerate(row):
                        built.cell(row_index, cell_index).text = value
            if comment is not None:
                document.add_comment(
                    runs=paragraph.runs, text=comment, author="Reviewer", initials="RV"
                )
    buffer = io.BytesIO()
    document.save(buffer)
    data = buffer.getvalue()
    if tracked:
        data = _rewrite_part(data, "word/document.xml", _insert_tracked_change)
    return data


def _png_bytes() -> bytes:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (64, 64), (10, 120, 200)).save(buffer, format="PNG")
    return buffer.getvalue()


def pptx_bytes(
    *,
    titles: Sequence[str] = ("Opening position", "The numbers", "What we decided"),
    body: str = "First point\nSecond point",
    notes: Sequence[str | None] = (None, "The number is soft; the trend is not.", None),
    picture_on: Sequence[int] = (),
) -> bytes:
    """A deck. `notes[i]` is slide i+1's speaker notes, or None for no notes.

    `picture_on` holds 1-indexed slide numbers that get an embedded PNG.
    """
    presentation = pptx.Presentation()
    picture = _png_bytes()
    for index, title in enumerate(titles):
        layout = presentation.slide_layouts[1 if index == 0 else 5]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        if index == 0 and body:
            slide.placeholders[1].text = body
        if index + 1 in picture_on:
            slide.shapes.add_picture(io.BytesIO(picture), Inches(1), Inches(2))
        note = notes[index] if index < len(notes) else None
        if note is not None:
            slide.notes_slide.notes_text_frame.text = note
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _cache_formula_values(xml: bytes) -> bytes:
    """Give every formula the cached `<v>` Excel writes and openpyxl does not.

    Without this, `data_only=True` returns None for every formula cell of any
    openpyxl-written workbook, and a test of the value path asserts nothing.
    """
    return re.sub(rb"(<f>[^<]*</f>)", rb"\1<v>%d</v>" % CACHED_FORMULA_VALUE, xml)


def xlsx_bytes(*, formulas: bool = False, cached: bool = False) -> bytes:
    """A two-sheet workbook.

    `formulas` adds a `=B*2` column; `cached` additionally gives those formulas
    the value Excel would have stored beside them.
    """
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["region", "units", "doubled"])
    for row, region in enumerate(("north", "south", "east"), start=2):
        sheet.append([region, row, None])
        if formulas:
            sheet[f"C{row}"] = f"=B{row}*2"
    notes = workbook.create_sheet("Notes")
    notes["A1"] = "Units are thousands."
    buffer = io.BytesIO()
    workbook.save(buffer)
    data = buffer.getvalue()
    if formulas and cached:
        data = _rewrite_part(data, "xl/worksheets/sheet1.xml", _cache_formula_values)
    return data


def big_xlsx(rows: int) -> bytes:
    """One wide sheet, for proving a budget truncates and says so."""
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Wide"
    sheet.append(["region", "units", "note"])
    for row in range(rows):
        sheet.append([f"region-{row}", row, "a note long enough to add up"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _odf_package(mimetype: str, content: str) -> bytes:
    """A three-entry ODF package, `mimetype` first and STORED.

    First and uncompressed is not a stylistic choice: the ODF specification
    requires it so the type can be sniffed from a file's opening bytes, and it
    is exactly what `adapters/ooxml.office_mimetype` relies on.
    """
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as package:
        stored = zipfile.ZipInfo("mimetype")
        stored.compress_type = zipfile.ZIP_STORED
        package.writestr(stored, mimetype)
        package.writestr("content.xml", content)
        package.writestr("META-INF/manifest.xml", _MANIFEST)
    return buffer.getvalue()


def odt_bytes(
    blocks: Sequence[tuple[str, str]] = (
        ("h", "Alpha"),
        ("p", "The body of the alpha section."),
        ("h", "Bravo"),
        ("p", "The body of the bravo section."),
    ),
) -> bytes:
    """An ODF text document. Each block is `("h", text)` or `("p", text)`."""
    body = "".join(
        f'<text:h text:outline-level="1">{text}</text:h>'
        if kind == "h"
        else f"<text:p>{text}</text:p>"
        for kind, text in blocks
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<office:document-content xmlns:office="{_OFFICE_NS}" xmlns:text="{_TEXT_NS}">'
        f"<office:body><office:text>{body}</office:text></office:body>"
        "</office:document-content>"
    )
    return _odf_package("application/vnd.oasis.opendocument.text", content)


def odp_bytes(
    slides: Sequence[tuple[str, Sequence[str]]] = (
        ("Opening position", ("First point", "Second point")),
        ("What we decided", ("Ship it",)),
    ),
) -> bytes:
    """An ODF presentation, one `draw:page` per slide."""
    pages = "".join(
        f'<draw:page draw:name="page{index}">'
        + f"<draw:frame><draw:text-box><text:p>{title}</text:p></draw:text-box></draw:frame>"
        + "".join(
            f"<draw:frame><draw:text-box><text:p>{line}</text:p></draw:text-box></draw:frame>"
            for line in lines
        )
        + "</draw:page>"
        for index, (title, lines) in enumerate(slides, start=1)
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<office:document-content xmlns:office="{_OFFICE_NS}" xmlns:text="{_TEXT_NS}" '
        f'xmlns:draw="{_DRAW_NS}">'
        f"<office:body><office:presentation>{pages}</office:presentation></office:body>"
        "</office:document-content>"
    )
    return _odf_package("application/vnd.oasis.opendocument.presentation", content)


def ods_bytes(
    sheets: Sequence[tuple[str, Sequence[Sequence[str]]]] = (
        ("Data", (("region", "units"), ("north", "2"), ("south", "3"))),
        ("Notes", (("Units are thousands.",),)),
    ),
) -> bytes:
    """An ODF spreadsheet, one `table:table` per sheet."""
    tables = "".join(
        f'<table:table table:name="{name}">'
        + "".join(
            "<table:table-row>"
            + "".join(
                f"<table:table-cell><text:p>{cell}</text:p></table:table-cell>" for cell in row
            )
            + "</table:table-row>"
            for row in rows
        )
        + "</table:table>"
        for name, rows in sheets
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<office:document-content xmlns:office="{_OFFICE_NS}" xmlns:text="{_TEXT_NS}" '
        f'xmlns:table="{_TABLE_NS}">'
        f"<office:body><office:spreadsheet>{tables}</office:spreadsheet></office:body>"
        "</office:document-content>"
    )
    return _odf_package("application/vnd.oasis.opendocument.spreadsheet", content)
