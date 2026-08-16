"""The fixtures are guarded, because a broken fixture makes a test vacuous.

Mirrors `tests/unit/test_pdf_fixtures.py`: if `cached=True` ever stops
injecting values, every formula-versus-value test silently becomes a test that
`None == None`.
"""

from __future__ import annotations

import io
import zipfile

import docx
import openpyxl
import pptx

from tests.fixtures_office import (
    CACHED_FORMULA_VALUE,
    big_xlsx,
    docx_bytes,
    odp_bytes,
    ods_bytes,
    odt_bytes,
    pptx_bytes,
    xlsx_bytes,
)


def test_the_cached_workbook_has_values_and_the_uncached_one_does_not() -> None:
    """openpyxl computes nothing and writes no cached value, so without the
    injection `data_only=True` returns None for every formula cell — and a test
    of the value path would be asserting that a blank is a blank.
    """
    bare = openpyxl.load_workbook(io.BytesIO(xlsx_bytes(formulas=True)), data_only=True)
    assert [c.value for c in bare["Data"][2]][2] is None

    cached = openpyxl.load_workbook(
        io.BytesIO(xlsx_bytes(formulas=True, cached=True)), data_only=True
    )
    assert [c.value for c in cached["Data"][2]][2] == CACHED_FORMULA_VALUE


def test_the_formula_workbook_still_reads_back_as_formulas() -> None:
    """Injecting a cached value must not destroy the `<f>` beside it."""
    book = openpyxl.load_workbook(io.BytesIO(xlsx_bytes(formulas=True, cached=True)))
    assert [c.value for c in book["Data"][2]][2] == "=B2*2"


def test_the_tracked_fixture_carries_a_tracked_change_and_the_plain_one_does_not() -> None:
    plain = zipfile.ZipFile(io.BytesIO(docx_bytes())).read("word/document.xml")
    tracked = zipfile.ZipFile(io.BytesIO(docx_bytes(tracked=True))).read("word/document.xml")
    assert b"<w:ins " not in plain
    assert b"<w:ins " in tracked


def test_the_tracked_fixture_is_still_a_readable_document() -> None:
    """A rewritten part python-docx cannot parse would make every tracked
    changes test fail for the wrong reason.
    """
    document = docx.Document(io.BytesIO(docx_bytes(tracked=True)))
    assert any("alpha" in p.text.lower() for p in document.paragraphs)


def test_the_commented_fixture_carries_exactly_one_comment() -> None:
    document = docx.Document(io.BytesIO(docx_bytes(comment="Check this number.")))
    comments = list(document.comments)
    assert len(comments) == 1
    assert comments[0].text == "Check this number."


def test_the_table_sits_between_headings_not_at_the_end() -> None:
    """A handler that appends tables rather than reading document order must
    produce visibly wrong text, so the fixture must not put the table last.
    """
    from docx.oxml.ns import qn

    document = docx.Document(io.BytesIO(docx_bytes()))
    tags = [child.tag for child in document.element.body.iterchildren()]
    assert qn("w:tbl") in tags
    assert tags.index(qn("w:tbl")) < len(tags) - 2


def test_notes_are_present_on_exactly_the_slides_that_asked_for_them() -> None:
    deck = pptx.Presentation(io.BytesIO(pptx_bytes()))
    assert [s.has_notes_slide for s in deck.slides] == [False, True, False]


def test_a_picture_lands_only_on_the_requested_slide() -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = pptx.Presentation(io.BytesIO(pptx_bytes(picture_on=(2,))))
    pictures = [
        sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        for slide in deck.slides
    ]
    assert pictures == [0, 1, 0]


def test_every_odf_package_stores_its_mimetype_first_and_uncompressed() -> None:
    """The ODF specification requires it and `office_mimetype` depends on it. A
    fixture that compressed the entry would make the ODF detection tests pass
    or fail for reasons unrelated to detection.
    """
    for data in (odt_bytes(), odp_bytes(), ods_bytes()):
        package = zipfile.ZipFile(io.BytesIO(data))
        assert package.namelist()[0] == "mimetype"
        assert package.getinfo("mimetype").compress_type == zipfile.ZIP_STORED


def test_a_big_workbook_is_actually_big() -> None:
    book = openpyxl.load_workbook(io.BytesIO(big_xlsx(500)), read_only=True)
    assert book["Wide"].max_row == 501
    book.close()
