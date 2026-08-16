"""Slide decks, OOXML and ODF.

The second producer of `PageRef`, and a better one than the first. A PDF's
pagination is typographic — page 7 is where the text happened to break. A
deck's is semantic: slide 7 is a thing an author made, so a barrier at every
slide boundary is the most natural barrier in any format this library reads.

Speaker notes are included in `represent` and LABELLED. They routinely hold the
reasoning the slide only asserts, so dropping them loses the argument and
keeping them unlabelled lets a model attribute a presenter's aside to the slide
itself. `NOTES_HEADING` is what separates the two claims.

Embedded pictures are reachable but never inlined: `list_media` reports what is
there and `describe_slide_image` hands one to the injected vision model. A
listing that returned megabytes of PNG would defeat the progressive disclosure
the card path exists for.

python-pptx is imported directly here, guarded exactly as `pdf.py` guards
pypdfium2.
"""

from __future__ import annotations

import io
import time
from dataclasses import dataclass
from typing import ClassVar

try:
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:  # pragma: no cover - exercised via a patched sys.modules
    raise ImportError(
        "readeverything's slide support needs python-pptx, which ships in the "
        "'office' extra: pip install 'readeverything[office]'. "
        "The composition root omits slide handling when python-pptx is absent, so "
        "reaching this means the handler was imported directly."
    ) from exc
from pydantic import BaseModel, Field

from readeverything.adapters.odf import odf_slides
from readeverything.adapters.ooxml import ODF_SLIDES_MIME, SLIDES_MIME, office_mimetype
from readeverything.domain.affordance import Affordance, DetailLevel
from readeverything.domain.capability import Capability
from readeverything.domain.card import Card, Segment
from readeverything.domain.errors import UnknownAffordanceError
from readeverything.domain.identity import MediaKind, SourceRef
from readeverything.domain.locator_map import LocatorMap, LocatorSegment
from readeverything.domain.locators import ByteRange, CharSpan, PageRef
from readeverything.domain.observation import OperationFinished, OperationStarted
from readeverything.domain.rendition import (
    Budget,
    Degradation,
    ImageContent,
    Rendered,
    Rendition,
    TextContent,
)
from readeverything.handlers.page_images import (
    PageImageParams,
    page_image_affordance,
    render_page_image,
)
from readeverything.ports.observation import Observer, emit
from readeverything.ports.rendering import DocumentRenderer
from readeverything.ports.source import SourceReader
from readeverything.ports.vision import VisionModel

#: What `page_image` renders one of. A deck's pagination is semantic --
#: slide 7 is a thing an author made -- so this is the format where a
#: rendering answers the most: arrangement, emphasis and imagery that no
#: text extraction recovers.
_RENDER_UNIT = "slide"

#: What `represent` calls itself when it narrates, matching every other handler.
_OPERATION = "represent"

#: What separates a presenter's aside from what the slide itself claims. A note
#: folded into the body without this reads as something the author put on
#: screen, which is a different assertion with a different audience.
NOTES_HEADING = "Speaker notes:"

#: Every slide's text ends with this, and the slide's `LocatorSegment` INCLUDES
#: it. `LocatorMap` demands total, gapless, zero-start coverage and
#: `CharSpan.__post_init__` rejects `start >= end`, so a slide holding only a
#: picture would otherwise contribute a zero-width span and raise. Owning the
#: separator means every slide owns at least one character no matter what it
#: contained. Do not "simplify" this away: an empty slide in the middle of a
#: deck is what breaks the map.
SLIDE_SEPARATOR = "\n"

#: What stands in the flattened text for a slide that carried no text at all. A
#: slide holding one diagram is not an empty slide, and saying nothing about it
#: would report the deck as shorter than it is.
_EMPTY_SLIDE = "(this slide carries no text)"


@dataclass(frozen=True, slots=True)
class _Picture:
    """One embedded image, located by the slide and position that reach it."""

    page: int
    index: int
    mime: str
    size_bytes: int


@dataclass(frozen=True, slots=True)
class _Slide:
    """One slide's title, body runs and notes, before any text is joined."""

    title: str
    body: tuple[str, ...]
    notes: str | None

    def rendered(self) -> str:
        parts = [part for part in (self.title, *self.body) if part]
        if self.notes:
            parts.append(f"{NOTES_HEADING} {self.notes}")
        return "\n".join(parts) if parts else _EMPTY_SLIDE


class ReadSlideParams(BaseModel):
    page: int = Field(default=1, ge=1, description="1-indexed slide number to read.")


class ListMediaParams(BaseModel):
    pass


class DescribeSlideImageParams(BaseModel):
    question: str = Field(description="What you want to know about the picture.")
    page: int = Field(default=1, ge=1, description="1-indexed slide number.")
    index: int = Field(default=0, ge=0, description="0-indexed picture number within the slide.")


class DescribeSlideParams(BaseModel):
    question: str = Field(description="What you want to know about the slide.")
    page: int = Field(default=1, ge=1, description="1-indexed slide number.")
    dpi: int = Field(default=150, gt=0, description="Render resolution, in dots per inch.")


class OfficeSlidesHandler:
    """Reads a deck, and maps every character to the slide it came from."""

    mime_patterns: ClassVar[tuple[str, ...]] = (SLIDES_MIME, ODF_SLIDES_MIME)
    priority: ClassVar[int] = 0
    handler_id: ClassVar[str] = "office_slides"
    handler_version: ClassVar[int] = 1

    def __init__(
        self,
        *,
        source: SourceReader,
        vision: VisionModel | None = None,
        renderer: DocumentRenderer | None = None,
        observer: Observer | None = None,
    ) -> None:
        self._source = source
        self._vision = vision
        self._renderer = renderer
        self._observer = observer

    def requires(self) -> frozenset[Capability]:
        """Nothing. Reading a deck's text needs no model and no binary."""
        return frozenset()

    def affordances(self) -> tuple[Affordance, ...]:
        affordances: list[Affordance] = [
            Affordance(
                name="read_slide",
                description="Return one slide's title, body text and speaker notes.",
                params=ReadSlideParams,
                requires=frozenset(),
                level=DetailLevel.SEGMENT,
            ),
            Affordance(
                name="list_media",
                description="List every picture embedded in the deck, with the slide it sits on.",
                params=ListMediaParams,
                requires=frozenset(),
                level=DetailLevel.SEGMENT,
            ),
        ]
        if self._vision is not None:
            affordances.append(
                Affordance(
                    name="describe_slide_image",
                    description=(
                        "Ask a vision model a question about a picture embedded in a "
                        "slide. Use this when the answer is in a chart or a diagram "
                        "rather than in the slide's own text — read_slide is far "
                        "cheaper when the text is really there."
                    ),
                    params=DescribeSlideImageParams,
                    requires=frozenset({Capability.VISION}),
                    level=DetailLevel.DEEP,
                )
            )
        if self._renderer is not None:
            # Declared, not published. The registry filters it out unless
            # DOCUMENT_RENDER is genuinely available, which is what keeps a
            # composition free to wire the converter unconditionally -- exactly
            # as it wires ffmpeg -- without an agent ever seeing a tool that
            # exists and apologises.
            #
            # Independent of `vision` above, and deliberately: a converter with
            # no vision model still hands back a slide image for something else
            # to read, and a vision model with no converter still describes an
            # embedded picture. Two capabilities, two affordances.
            affordances.append(page_image_affordance(_RENDER_UNIT))
        if self._renderer is not None and self._vision is not None:
            # Two capabilities at once, which is what earns this its own
            # affordance rather than leaving a caller to chain `page_image`
            # into `ask_about_image`. That chain costs two round trips and a
            # schema read to express one intent -- the argument that produced
            # `ask_about_image` itself, one medium along.
            #
            # `describe_slide_image` above is a DIFFERENT question and both
            # exist deliberately: that one asks about a picture the author
            # embedded, this one asks about the slide as the audience saw it.
            affordances.append(
                Affordance(
                    name="describe_slide",
                    description=(
                        "Render a whole slide and ask a vision model about it. Use "
                        "this when the answer is in the slide's arrangement, a chart "
                        "or a diagram rather than in its text -- read_slide is far "
                        "cheaper when the text is really there. The model sees a "
                        "faithful rendering of the slide, not the original file."
                    ),
                    params=DescribeSlideParams,
                    requires=frozenset({Capability.DOCUMENT_RENDER, Capability.VISION}),
                    level=DetailLevel.DEEP,
                )
            )
        return tuple(affordances)

    # -- parsing ---------------------------------------------------------

    def _is_odf(self, data: bytes) -> bool:
        return office_mimetype(data) == ODF_SLIDES_MIME

    def _slides_from_odf(self, data: bytes) -> list[_Slide] | None:
        pages = odf_slides(data)
        if not pages:
            return None
        # The first run of a `draw:page` is its title by convention, which is
        # both how ODP stores it and how a reader reads it. ODP carries no
        # speaker notes through this reader, so `notes` is honestly None rather
        # than an empty string that would render a bare heading.
        return [
            _Slide(title=runs[0] if runs else "", body=tuple(runs[1:]), notes=None)
            for runs in pages
        ]

    def _slides_from_ooxml(self, data: bytes) -> list[_Slide] | None:
        try:
            presentation = pptx.Presentation(io.BytesIO(data))
            slides: list[_Slide] = []
            for slide in presentation.slides:
                title_shape = slide.shapes.title
                title = " ".join(title_shape.text.split()) if title_shape is not None else ""
                body: list[str] = []
                for shape in slide.shapes:
                    if shape is title_shape or not shape.has_text_frame:
                        continue
                    text = "\n".join(
                        line for line in shape.text_frame.text.splitlines() if line.strip()
                    )
                    if text:
                        body.append(text)
                slides.append(_Slide(title=title, body=tuple(body), notes=self._notes_of(slide)))
        except Exception:
            return None
        return slides

    def _notes_of(self, slide: object) -> str | None:
        """A slide's speaker notes, or None when it has none.

        `has_notes_slide` is checked FIRST and is not optional: reading
        `slide.notes_slide` CREATES a notes slide as a side effect, so an
        unguarded read mutates the parsed package and makes the card's
        `notes_present` fact depend on whether anything looked at the deck
        earlier. That bug only ever shows on the second call.
        """
        if not getattr(slide, "has_notes_slide", False):
            return None
        text = slide.notes_slide.notes_text_frame.text  # type: ignore[attr-defined]
        collapsed = " ".join(text.split())
        return collapsed or None

    def _parse(self, data: bytes) -> list[_Slide] | None:
        if self._is_odf(data):
            return self._slides_from_odf(data)
        return self._slides_from_ooxml(data)

    def _pictures(self, data: bytes) -> tuple[_Picture, ...]:
        """Every embedded picture's shape, WITHOUT decoding its bytes.

        `shape.image.blob` is deliberately not touched here: the card path and
        `list_media` both want the inventory, and materialising every PNG to
        count them would make a listing cost what a download costs.
        """
        if self._is_odf(data):
            return ()
        found: list[_Picture] = []
        try:
            presentation = pptx.Presentation(io.BytesIO(data))
            for number, slide in enumerate(presentation.slides, start=1):
                index = 0
                for shape in slide.shapes:
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                    image = shape.image
                    found.append(
                        _Picture(
                            page=number,
                            index=index,
                            mime=image.content_type,
                            size_bytes=len(image.blob),
                        )
                    )
                    index += 1
        except Exception:
            return ()
        return tuple(found)

    def _picture_blob(self, data: bytes, page: int, index: int) -> tuple[bytes, str] | None:
        if self._is_odf(data):
            return None
        try:
            presentation = pptx.Presentation(io.BytesIO(data))
            slides = list(presentation.slides)
            if page > len(slides):
                return None
            found = 0
            for shape in slides[page - 1].shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                if found == index:
                    return shape.image.blob, shape.image.content_type
                found += 1
        except Exception:
            return None
        return None

    # -- flattening ------------------------------------------------------

    def _flatten(
        self, slides: list[_Slide]
    ) -> tuple[str, tuple[LocatorSegment, ...], tuple[int, ...]]:
        chunks: list[str] = []
        segments: list[LocatorSegment] = []
        barriers: list[int] = []
        cursor = 0
        for index, slide in enumerate(slides):
            number = index + 1  # `PageRef` counts as a reader does; pptx does not.
            chunk = slide.rendered() + SLIDE_SEPARATOR
            if index:
                # A barrier is where a new slide's first character begins, so
                # there is exactly one per slide break: slide count minus one.
                barriers.append(cursor)
            segments.append(LocatorSegment(CharSpan(cursor, cursor + len(chunk)), PageRef(number)))
            cursor += len(chunk)
            chunks.append(chunk)
        return "".join(chunks), tuple(segments), tuple(barriers)

    # -- the handler surface ---------------------------------------------

    async def describe(self, ref: SourceRef) -> Card:
        """Slide count, notes and media inventory. No vision model is touched.

        `kind` is `BINARY`, matching `pdf.py`: these mimetypes reach this
        handler at the registry's exact-mimetype step, long before the kind
        step, so what identifies a deck's card is its mime, its slide count and
        its affordances.
        """
        data = await self._source.read_bytes(ref.uri)
        slides = self._parse(data)
        if slides is None:
            return Card(
                ref=ref,
                kind=MediaKind.BINARY,
                facts={"readable": "no", "size_bytes": ref.size_bytes},
                outline=(),
                excerpt=None,
                affordances=self.affordances(),
            )
        pictures = self._pictures(data)
        return Card(
            ref=ref,
            kind=MediaKind.BINARY,
            facts={
                "readable": "yes",
                "slide_count": len(slides),
                "notes_present": "yes" if any(s.notes for s in slides) else "no",
                "media_count": len(pictures),
                "size_bytes": ref.size_bytes,
            },
            outline=tuple(
                Segment(PageRef(index + 1), slide.title or f"slide {index + 1}")
                for index, slide in enumerate(slides)
            ),
            excerpt=None,
            affordances=self.affordances(),
        )

    async def invoke(self, ref: SourceRef, name: str, params: BaseModel) -> Rendition:
        match name:
            case "read_slide":
                if not isinstance(params, ReadSlideParams):
                    raise TypeError(f"expected ReadSlideParams, got {type(params).__name__}")
                return await self._read_slide(ref, params.page)
            case "list_media":
                if not isinstance(params, ListMediaParams):
                    raise TypeError(f"expected ListMediaParams, got {type(params).__name__}")
                return await self._list_media(ref)
            case "describe_slide_image":
                if not isinstance(params, DescribeSlideImageParams):
                    raise TypeError(
                        f"expected DescribeSlideImageParams, got {type(params).__name__}"
                    )
                return await self._describe_slide_image(ref, params)
            case "page_image":
                if not isinstance(params, PageImageParams):
                    raise TypeError(f"expected PageImageParams, got {type(params).__name__}")
                return await self._page_image(ref, params)
            case "describe_slide":
                if not isinstance(params, DescribeSlideParams):
                    raise TypeError(f"expected DescribeSlideParams, got {type(params).__name__}")
                return await self._describe_slide(ref, params)
            case _:
                raise UnknownAffordanceError(name, (a.name for a in self.affordances()))

    async def _page_image(self, ref: SourceRef, params: PageImageParams) -> Rendition:
        if self._renderer is None:
            raise UnknownAffordanceError("page_image", (a.name for a in self.affordances()))
        return await render_page_image(
            renderer=self._renderer,
            source=self._source,
            ref=ref,
            params=params,
            unit=_RENDER_UNIT,
        )

    def _degraded(self, ref: SourceRef, detail: str) -> Rendition:
        """What every unreadable or out-of-range request returns.

        Never an exception: an agent guessing a slide number gets a result it
        can read and correct.
        """
        return Rendition(
            locator=ByteRange(0, max(1, ref.size_bytes)),
            content=TextContent(detail),
            degraded=True,
        )

    async def _read_slide(self, ref: SourceRef, page: int) -> Rendition:
        data = await self._source.read_bytes(ref.uri)
        slides = self._parse(data)
        if slides is None:
            return self._degraded(ref, f"{ref.uri} could not be opened as a slide deck")
        if page > len(slides):
            return self._degraded(
                ref, f"slide {page} does not exist; the deck has {len(slides)} slide(s)"
            )
        return Rendition(locator=PageRef(page), content=TextContent(slides[page - 1].rendered()))

    async def _list_media(self, ref: SourceRef) -> Rendition:
        data = await self._source.read_bytes(ref.uri)
        # Readability is checked FIRST, and the order is the whole point: a
        # file that could not be opened embeds no pictures AND nothing else, so
        # reporting "embeds no pictures" about it is a claim nothing
        # established. "There are none" must mean the deck was read.
        if self._parse(data) is None:
            return self._degraded(ref, f"{ref.uri} could not be opened as a slide deck")
        pictures = self._pictures(data)
        if not pictures:
            # "There are none" and empty output are different answers, and only
            # one of them tells an agent to stop looking.
            return Rendition(
                locator=ByteRange(0, max(1, ref.size_bytes)),
                content=TextContent("the deck embeds no pictures"),
            )
        lines = [
            f"slide {p.page}, picture {p.index}: {p.mime}, {p.size_bytes} bytes" for p in pictures
        ]
        return Rendition(
            locator=ByteRange(0, max(1, ref.size_bytes)),
            content=TextContent("\n".join(lines)),
        )

    async def _describe_slide_image(
        self, ref: SourceRef, params: DescribeSlideImageParams
    ) -> Rendition:
        if self._vision is None:
            raise UnknownAffordanceError(
                "describe_slide_image", (a.name for a in self.affordances())
            )
        data = await self._source.read_bytes(ref.uri)
        found = self._picture_blob(data, params.page, params.index)
        if found is None:
            return self._degraded(
                ref, f"slide {params.page} has no picture at index {params.index}"
            )
        blob, mime = found
        try:
            answer = await self._vision.describe(blob, mime, params.question)
        except Exception:
            return Rendition(
                locator=PageRef(params.page),
                content=TextContent(
                    f"the vision model failed to answer about the picture on slide {params.page}"
                ),
                degraded=True,
            )
        if not answer.strip():
            return Rendition(
                locator=PageRef(params.page),
                content=TextContent(
                    f"the vision model returned no answer about the picture on slide {params.page}"
                ),
                degraded=True,
            )
        return Rendition(
            locator=PageRef(params.page),
            content=TextContent(" ".join(answer.split())),
        )

    async def _describe_slide(self, ref: SourceRef, params: DescribeSlideParams) -> Rendition:
        if self._vision is None or self._renderer is None:
            raise UnknownAffordanceError("describe_slide", (a.name for a in self.affordances()))
        # Rendering first, and its degradation is returned UNCHANGED when it
        # fails: the reason a slide could not be rendered is a better answer
        # than "the model could not describe it", and it is the true one.
        rendered = await render_page_image(
            renderer=self._renderer,
            source=self._source,
            ref=ref,
            params=PageImageParams(page=params.page, dpi=params.dpi),
            unit=_RENDER_UNIT,
        )
        if not isinstance(rendered.content, ImageContent):
            return rendered
        try:
            answer = await self._vision.describe(
                rendered.content.data, rendered.content.mime, params.question
            )
        except Exception:
            return self._degraded(
                ref, f"the vision model failed to answer about slide {params.page}"
            )
        if not answer.strip():
            return self._degraded(
                ref, f"the vision model returned no answer about slide {params.page}"
            )
        return Rendition(
            locator=PageRef(params.page),
            content=TextContent(" ".join(answer.split())),
            # The model looked at a RENDERING, so an answer about typography is
            # an answer about the converter's font substitutions. Carrying the
            # provenance through is what stops that being read as a fact about
            # the deck.
            degradations=rendered.degradations,
        )

    async def represent(self, ref: SourceRef, budget: Budget) -> Rendered:
        """Narrated start to finish, matching every other handler."""
        emit(self._observer, OperationStarted(operation=_OPERATION, ref=ref))
        started = time.perf_counter()
        try:
            return await self._represent(ref, budget)
        finally:
            emit(
                self._observer,
                OperationFinished(
                    operation=_OPERATION, ref=ref, elapsed_s=time.perf_counter() - started
                ),
            )

    async def _represent(self, ref: SourceRef, budget: Budget) -> Rendered:
        data = await self._source.read_bytes(ref.uri)
        slides = self._parse(data)
        if slides is None:
            return self._nothing_to_read(
                ref,
                budget,
                summary=f"Unreadable slide deck {ref.uri}, {ref.size_bytes} bytes.",
                what="deck unopenable",
                detail="the file could not be opened as a slide deck; no text was extracted",
            )
        if not slides:
            return self._nothing_to_read(
                ref,
                budget,
                summary=f"Slide deck {ref.uri} has no slides.",
                what="deck has no slides",
                detail="the deck opened but contains no slides; no text was extracted",
            )
        text, segments, barriers = self._flatten(slides)
        return self._fit(text, segments, barriers, budget, ())

    def _nothing_to_read(
        self, ref: SourceRef, budget: Budget, *, summary: str, what: str, detail: str
    ) -> Rendered:
        """A rendition for a file with no slide to point at.

        Located by `ByteRange` rather than `PageRef`: no slide was ever
        observed, and claiming slide 1 would be a claim about a deck this
        handler never read a slide of.
        """
        segments = (
            LocatorSegment(CharSpan(0, len(summary)), ByteRange(0, max(1, ref.size_bytes))),
        )
        return self._fit(summary, segments, (), budget, (Degradation(what=what, detail=detail),))

    def _fit(
        self,
        full: str,
        segments: tuple[LocatorSegment, ...],
        barriers: tuple[int, ...],
        budget: Budget,
        degradations: tuple[Degradation, ...],
    ) -> Rendered:
        """Apply the budget, pruning the map and the barriers along with the text.

        `Rendered` rejects a map that does not cover its text exactly and a
        barrier past the end, so truncation cannot touch the text alone.
        """
        if budget.max_chars is None or len(full) <= budget.max_chars:
            return Rendered(
                text=full,
                locator_map=LocatorMap.build(segments),
                barriers=barriers,
                degradations=degradations,
            )
        keep = max(1, budget.max_chars)
        text = full[:keep]
        kept = tuple(
            LocatorSegment(CharSpan(s.span.start, min(s.span.end, keep)), s.locator)
            for s in segments
            if s.span.start < keep
        )
        return Rendered(
            text=text,
            locator_map=LocatorMap.build(kept),
            barriers=tuple(barrier for barrier in barriers if barrier < keep),
            degradations=(
                *degradations,
                Degradation(
                    what="text truncated",
                    detail=f"kept {len(text)} of {len(full)} characters",
                ),
            ),
        )
